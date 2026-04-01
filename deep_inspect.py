import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

prs = Presentation(r"C:\Doc Ref\MS24147_Final_Presentation.pptx")

print(f"Total slides: {len(prs.slides)}")
print(f"Slide Width: {prs.slide_width.inches:.3f} in, Height: {prs.slide_height.inches:.3f} in")

for slide_idx in range(min(5, len(prs.slides))):
    slide = prs.slides[slide_idx]
    print(f"\n====== SLIDE {slide_idx+1} ======")

    # Background
    fill = slide.background.fill
    print(f"  BG fill type: {fill.type}")
    try:
        print(f"  BG solid color: #{fill.fore_color.rgb}")
    except:
        print("  BG color: from master")

    # Shapes
    for shape in slide.shapes:
        print(f"\n  [Shape] Name={shape.name} | Type={shape.shape_type} | Left={shape.left} Top={shape.top} W={shape.width} H={shape.height}")
        
        # Fill
        try:
            sf = shape.fill
            print(f"    Fill type: {sf.type}")
            if sf.type == 1:
                print(f"    Fill color: #{sf.fore_color.rgb}")
        except: pass
        
        # Line
        try:
            ln = shape.line
            if ln.color.type:
                print(f"    Line color: #{ln.color.rgb}, Width: {ln.width}")
        except: pass

        # Text  
        if shape.has_text_frame:
            for pi, para in enumerate(shape.text_frame.paragraphs):
                for run in para.runs:
                    try:
                        col = f"#{run.font.color.rgb}" if run.font.color.type else "Auto"
                    except:
                        col = "N/A"
                    print(f"    Para[{pi}] Run: '{run.text[:50]}' | Font={run.font.name} Sz={run.font.size} Bold={run.font.bold} Color={col}")
