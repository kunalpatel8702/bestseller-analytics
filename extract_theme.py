from pptx import Presentation

try:
    prs = Presentation(r"C:\Doc Ref\MS24147_Final_Presentation.pptx")
    slide = prs.slides[0]
    
    print("--- Theme Analysis ---")
    if slide.background.fill.type == 1: # solid
        print(f"Background Solid Color: {slide.background.fill.fore_color.rgb}")
    elif slide.background.fill.type is None:
        print("Background is inherited from master.")
        
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    print(f"Text: '{run.text}', Font Name: {run.font.name}, Size: {run.font.size}, Color: {run.font.color.rgb if run.font.color.type else 'Auto'}")
except Exception as e:
    print(f"Error: {e}")
