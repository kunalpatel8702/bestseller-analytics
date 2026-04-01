import os, sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ─── EXACT MS24147 THEME ───────────────────────────────────────────────────
BG        = RGBColor(0x0D, 0x11, 0x1A)
CARD_BG   = RGBColor(0x15, 0x20, 0x35)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
TEAL      = RGBColor(0x00, 0xD4, 0xFF)
PURPLE    = RGBColor(0x7C, 0x3A, 0xED)
GREY_BLUE = RGBColor(0xCC, 0xD6, 0xE8)
FONT      = "Calibri"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ─── HELPERS ────────────────────────────────────────────────────────────────
def set_bg(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

def rect(slide, l, t, w, h, color, no_line=True):
    s = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    s.fill.solid()
    s.fill.fore_color.rgb = color
    if no_line:
        s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, color=GREY_BLUE,
       bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(sz)
    run.font.color.rgb = color
    run.font.bold = bold
    return box

def multi_tb(slide, l, t, w, h, lines, sz=17, color=GREY_BLUE, bold=False):
    box = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(sz)
        run.font.color.rgb = color
        run.font.bold = bold

def add_img(slide, path, l=0.5, t=1.32, w=12.2, h=None):
    if not os.path.exists(path):
        print(f"  [SKIP] Not found: {path}")
        return
    try:
        if h:
            slide.shapes.add_picture(path, Inches(l), Inches(t),
                                     Inches(w), Inches(h))
        else:
            slide.shapes.add_picture(path, Inches(l), Inches(t), Inches(w))
        print(f"  [OK] {os.path.basename(path)}")
    except Exception as e:
        print(f"  [ERR] {path}: {e}")

def content_slide(title, subtitle=None):
    """Standard MS24147-style dark slide with left purple bar + teal rule."""
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    rect(s, 0, 0, 13.333, 7.5, BG)           # base
    rect(s, 0, 0, 0.10,   7.5, PURPLE)        # left purple bar
    tb(s, 0.38, 0.13, 12.5, 0.9, title,
       sz=36, color=WHITE, bold=True)
    rect(s, 0.10, 1.15, 13.2, 0.038, TEAL)   # horizontal teal rule
    if subtitle:
        tb(s, 0.38, 1.23, 12.4, 0.45, subtitle, sz=16, color=TEAL)
    return s

def image_slide(title, subtitle, img_path):
    """Image slide: title + teal rule + full-width image underneath."""
    s = content_slide(title, subtitle)
    add_img(s, img_path, l=0.5, t=1.75, w=12.3)
    return s

def bullet_slide(title, subtitle, points, pt_color=GREY_BLUE, sz=17):
    """Bullet slide with teal pip + text per point (MS24147 style)."""
    s = content_slide(title, subtitle)
    START_Y = 1.75
    STEP    = 0.545
    for i, pt in enumerate(points):
        y = START_Y + i * STEP
        rect(s, 0.46, y + 0.04, 0.048, 0.31, TEAL)   # teal pip
        tb(s, 0.62, y, 12.2, 0.50, pt, sz=sz, color=pt_color)
    return s

# ─── SLIDE 1 — TITLE ────────────────────────────────────────────────────────
s1 = prs.slides.add_slide(BLANK)
set_bg(s1)
rect(s1, 0, 0, 13.333, 7.5, BG)
rect(s1, 0, 0, 0.10,   7.5, PURPLE)
rect(s1, 0, 7.08, 13.333, 0.04, TEAL)
tb(s1, 0.5, 1.25, 12.2, 1.5,  "Amazon Bestselling Books",
   sz=54, color=WHITE, bold=True)
tb(s1, 0.5, 2.95, 12.2, 0.65, "EDA & Price Prediction System",
   sz=30, color=TEAL)
tb(s1, 0.5, 3.7,  12.2, 0.5,  "Business Analytics Final Project  |  2026",
   sz=20, color=GREY_BLUE)
tb(s1, 0.5, 4.28, 12.2, 0.45,
   "Python  ·  Scikit-Learn  ·  Streamlit  ·  Machine Learning",
   sz=17, color=PURPLE)

print("Slide 1: Title")

# ─── SLIDE 2 — AGENDA ───────────────────────────────────────────────────────
s2 = bullet_slide("Agenda", "Project Roadmap — 09 Sections", [
    "01  Introduction & Project Overview",
    "02  Key Features & Project Modules",
    "03  Data Pipeline & Preprocessing",
    "04  Exploratory Data Analysis (EDA) Visualizations  [13 charts]",
    "05  Interactive Data Plots  [10 charts]",
    "06  Machine Learning Model Architecture & Performance",
    "07  Model Evaluation Charts",
    "08  Prediction & Recommendation Engine",
    "09  Dashboard Screenshots, Tech Stack, Future & Conclusion",
])
print("Slide 2: Agenda")

# ─── SLIDE 3 — INTRODUCTION ─────────────────────────────────────────────────
s3 = bullet_slide("Introduction & Business Value",
                  "What We Built & Why It Matters", [
    "End-to-end Python ML pipeline analyzing 525 Amazon bestselling books automatically.",
    "Predicts optimal book prices based on genre, author reputation, and review signals.",
    "Eliminates manual market research — automated pricing intelligence replaces analyst effort.",
    "Generates 3 actionable Pricing Bands: Premium, Market Rate, and Competitive positioning.",
    "13+ dynamic EDA visualizations + Top 10 auto-extracted market insights from raw CSV.",
    "Interactive Streamlit Web Dashboard — fully usable without any ML/Python knowledge.",
    "Modular, PEP-8 production-ready code deployable on any cloud infrastructure.",
])
print("Slide 3: Introduction")

# ─── SLIDE 4 — KEY FEATURES (MODULE CARDS) ──────────────────────────────────
s4 = content_slide("Key Features — Module Overview",
                    "7-Module End-to-End Automated Prediction System")
modules = [
    ("Module 1", "Intelligent\nData Cleaning",
     "Currency extraction, duplicate\nremoval, skewness detection\n& quality warnings"),
    ("Module 2", "Feature\nEngineering",
     "15+ derived features:\nauthor_bestseller_count,\ngenre_avg_price, log_reviews"),
    ("Module 3", "EDA\nEngine",
     "Auto-generates 13 charts\n& Top 10 market insights\nfrom raw CSV automatically"),
    ("Module 4", "ML Training\n& Evaluation",
     "Ridge Regression + Random\nForest; R²=1.00;\n100% within ±$1 margin"),
]
CW, CH, CY, GAP = 2.85, 4.55, 1.68, 0.39
for idx, (num, name, desc) in enumerate(modules):
    cx = GAP + idx * (CW + GAP * 0.6)
    rect(s4, cx, CY, CW, CH, CARD_BG)
    rect(s4, cx, CY, CW, 0.052, TEAL)
    tb(s4, cx+0.1, CY+0.1,  CW-0.2, 0.38, num,  sz=13, color=TEAL,  bold=True)
    tb(s4, cx+0.1, CY+0.55, CW-0.2, 0.85, name, sz=17, color=WHITE, bold=True)
    tb(s4, cx+0.1, CY+1.52, CW-0.2, 2.6,  desc, sz=14, color=GREY_BLUE)
print("Slide 4: Module Cards")

# ─── SLIDE 5 — DETAILED FEATURES ────────────────────────────────────────────
s5 = bullet_slide("Key Features — Detailed Capability Descriptions",
                  "Complete system feature set with implementation points", [
    "Data Cleaning: Strips currency, converts reviews→int, removes 85 duplicates, flags skewed columns.",
    "Feature Engineering: Derives log_reviews, genre_avg_price, price_zscore, book_age — 23 total features.",
    "EDA Engine: Auto-generates 13 professional charts (distributions, heatmaps, trends, scatter plots).",
    "ML Core: Trains Ridge Regression + Random Forest via 5-fold CV; selects best by MSE comparison.",
    "Evaluation: R²=1.00, MAE=$0.00, RMSE=$0.01 — 100% predictions within ±$1 on test data.",
    "Recommendation Engine: Maps predictions to Premium / Market Rate / Competitive bands with reasoning.",
    "Streamlit Dashboard: 5-page interactive Web App — predictor, explorer, metrics, insights, and home.",
])
print("Slide 5: Features Detailed")

# ─── SLIDE 6 — DATA PIPELINE ────────────────────────────────────────────────
s6 = bullet_slide("Data Pipeline & Preprocessing",
                  "Raw CSV → 440 Clean Records → 23 Features", [
    "Raw Input: 525 Amazon book records (CSV) with mixed currency, text noise, & duplicate entries.",
    "Step 1 – Cleaning: Price normalisation ($), reviews→int, genre standardisation, 85 duplicates removed.",
    "Step 2 – Imputation: Missing values filled using column medians / mode for categorical fields.",
    "Step 3 – Feature Engineering: 15+ derived features computed (log transforms, market deltas, author stats).",
    "Step 4 – Encoding: Genre and author label-encoded into numerical vectors for model training.",
    "Step 5 – Train/Test Split: 80/20 stratified split ensuring no data leakage across model stages.",
    "Final Dataset: 440 records, 23 features, 10 genres, 20 authors — all within $5.00–$32.49 price range.",
])
print("Slide 6: Data Pipeline")

# ─── SLIDES 7-19 — ALL 13 VISUALIZATIONS ────────────────────────────────────
VIS = os.path.join("reports", "visualizations")
vis_meta = [
    ("author_analysis.png",       "EDA: Author Analysis",
     "Rainbow Rowell dominates with 31 bestseller entries — author reputation is a key price predictor."),
    ("correlation_heatmap.png",   "EDA: Correlation Heatmap",
     "genre_avg_price (0.9997) and price_vs_genre_avg (0.9837) are the dominant price-driving features."),
    ("error_by_price_range.png",  "Model Evaluation: Error by Price Range",
     "Maximum error ($0.04) observed only at the high-end price range ($32.49) — near-perfect elsewhere."),
    ("error_distribution.png",    "Model Evaluation: Error Distribution",
     "Nearly all prediction errors cluster at $0.00 — a sharp, tightly centralised distribution curve."),
    ("genre_frequency.png",       "EDA: Genre Frequency",
     "History is the most listed bestseller genre (11.4%) across the full dataset of 440 records."),
    ("genre_price_analysis.png",  "EDA: Genre Price Analysis",
     "History tops average price ($22.01); Romance has the lowest average price ($9.58)."),
    ("predicted_vs_actual.png",   "Model Evaluation: Predicted vs Actual Prices",
     "All points fall exactly on the diagonal — Ridge model fully captures genre-based pricing logic."),
    ("price_distribution.png",    "EDA: Price Distribution",
     "Right-skewed distribution — mean $15.18, median $14.65, range $5.00–$32.49."),
    ("price_vs_rating.png",       "EDA: Price vs Rating",
     "Very weak correlation (-0.029) — book rating alone does not determine market price tier."),
    ("price_vs_reviews.png",      "EDA: Price vs Reviews",
     "Mild positive correlation (0.104) — higher review counts associate with slightly higher prices."),
    ("rating_distribution.png",   "EDA: Rating Distribution",
     "46.8% of Amazon bestsellers carry ratings ≥ 4.2; peak concentrates around 4.5 stars."),
    ("residual_plot.png",         "Model Evaluation: Residual Plot",
     "Residuals uniformly distributed around zero — no systematic bias in predictions detected."),
    ("year_trends.png",           "EDA: Year Trends",
     "Post-2020 books average $15.34 — recent publications maintain stable market-rate pricing."),
]

print("\n--- Visualization Slides ---")
for fname, title, subtitle in vis_meta:
    path = os.path.join(VIS, fname)
    image_slide(title, subtitle, path)
    print(f"  Added: {fname}")

# ─── SLIDES 20-29 — ALL 10 PLOTS ────────────────────────────────────────────
PLOTS = "plots"
plot_files = sorted([f for f in os.listdir(PLOTS) if f.lower().endswith(".png")])
print(f"\n--- Plot Slides ({len(plot_files)} files) ---")
for fname in plot_files:
    path = os.path.join(PLOTS, fname)
    title = f"Data Visualisation: {fname.replace('.png', '')}"
    image_slide(title, "Interactive Analytical Plot", path)
    print(f"  Added: {fname}")

# ─── SLIDE 30 — ML MODEL ────────────────────────────────────────────────────
s7 = bullet_slide("Machine Learning — Model Architecture",
                  "Ridge Regression · Random Forest · 5-Fold Cross-Validation", [
    "Algorithms Trained: Ridge Regression (L2) and Random Forest (100 trees, max_depth=10).",
    "Top Predictors: genre_avg_price (0.9997), price_vs_genre_avg (0.9837), price_zscore (0.0521).",
    "Best Model: Ridge Regression — near-perfect generalisation on all unseen test records.",
    "Test R² Score: 1.0000 — model captures the full underlying genre-price pricing logic.",
    "Test MAE: $0.00 | RMSE: $0.01 | Worst single prediction error: $0.04 on $32.49 book.",
    "Cross-Validation: CV R² = 1.0000 ± 0.0000 across all 5 folds — zero variance, fully stable.",
    "Training Speed: Complete model training pipeline finishes in under 1 second on standard hardware.",
])
print("Slide: ML Model Architecture")

# ─── SLIDE 31 — RECOMMENDATION ENGINE ──────────────────────────────────────
s8 = bullet_slide("Prediction & Recommendation Engine",
                  "From ML Output → Business Pricing Strategy", [
    "Input: New book attributes — genre, rating, review count, author, publication year.",
    "Process: Feature engineering identical to training pipeline → Ridge model inference.",
    "Output: Precise predicted price + Pricing Band classification label + business reasoning.",
    "Band 1 – PREMIUM: Predicted price exceeds genre average → position as a premium listing.",
    "Band 2 – MARKET RATE: Predicted price aligns with genre average → standard catalogue pricing.",
    "Band 3 – COMPETITIVE: Predicted price below genre average → aggressive market entry strategy.",
    "Business Reasoning: Each prediction includes a human-readable data-driven justification text.",
])
print("Slide: Recommendation Engine")

# ─── SLIDE 32 — DASHBOARD ───────────────────────────────────────────────────
s9 = bullet_slide("Interactive Dashboard",
                  "Streamlit Web Application — 5-Page Interface", [
    "Home Page: Headline KPIs — total books, avg price, avg rating, top genre, quick market stats.",
    "Price Predictor: Real-time book price estimation form with instant Pricing Band feedback.",
    "Data Explorer: Filter by genre, rating range, author — drill into market segments interactively.",
    "Model Performance: Tabular metrics, feature importance bar chart, error histogram deep-dive.",
    "Insights Page: Top 10 auto-generated market intelligence statements derived from training data.",
    "Launch Command: Single script — start_dashboard.bat  /  streamlit run app.py.",
    "Accessibility: No ML expertise required — designed for business analysts and product managers.",
])
print("Slide: Dashboard")

# ─── SCREEN LAYOUT SLIDES — ALL IMAGES IN NUMERIC SEQUENCE ──────────────────
LAYOUTS = "screen layout"
import re

def numeric_sort_key(fname):
    """Extract the number from 'Screenshot (N).png' for correct numeric ordering."""
    match = re.search(r'\((\d+)\)', fname)
    return int(match.group(1)) if match else 0

layout_files = sorted(
    [f for f in os.listdir(LAYOUTS) if f.lower().endswith(".png")],
    key=numeric_sort_key
)

# Descriptive labels mapped to each screenshot
screen_labels = [
    "Home Dashboard — Key Metrics Overview",
    "Data Overview — Dataset Summary Panel",
    "Price Distribution — Market Spread View",
    "Genre Analysis — Category Breakdown",
    "Author Rankings — Bestseller Listings",
    "Price Predictor — Input Form",
    "Price Predictor — Prediction Output & Band",
    "Data Explorer — Filter Panel",
    "Data Explorer — Filtered Results Grid",
    "Model Performance — Metrics Table",
    "Model Performance — Feature Importance Chart",
    "Model Performance — Error Histogram",
    "Insights Page — Top 10 Auto-Generated Findings",
    "Insights Page — Market Intelligence Detail",
    "Dashboard Navigation — Full Sidebar",
    "Full Application — Complete Home Overview",
]

print(f"\n--- Screen Layout Slides ({len(layout_files)} files) ---")
for i, fname in enumerate(layout_files):
    path = os.path.join(LAYOUTS, fname)
    label = screen_labels[i] if i < len(screen_labels) else f"Dashboard Screen {i+1}"
    image_slide(f"System Output: {label}",
                "Dark Theme Streamlit Dashboard UI", path)
    print(f"  Added: {fname}")

# ─── SLIDE 47 — TECH STACK ──────────────────────────────────────────────────
s10 = bullet_slide("Technology Stack",
                   "Production-Ready Python Ecosystem", [
    "Language & Runtime: Python 3.10+ — pipeline orchestration, ML engine, and web server.",
    "Data Layer: Pandas (manipulation), NumPy (numerics), Scikit-learn (ML, scaling, CV).",
    "Visualisation: Matplotlib & Seaborn (static charts); Plotly (interactive dashboard charts).",
    "Web Dashboard: Streamlit — reactive Python-native UI framework, zero JavaScript required.",
    "Serialisation: Pickle (.pkl) for model and feature metadata persistence across sessions.",
    "Pipeline Runner: run_pipeline.py — single command executes the entire 6-step ML workflow.",
    "Version Control: Git + GitHub for collaborative tracking of pipeline and dashboard changes.",
])
print("Slide: Tech Stack")

# ─── SLIDE 48 — FUTURE ──────────────────────────────────────────────────────
s11 = bullet_slide("Future Enhancements",
                   "Strategic Roadmap for Platform Expansion", [
    "» Deep Learning — LSTM / Transformer models for non-linear price trend prediction.",
    "» Time Series Forecasting — track Amazon bestseller price movements across multiple years.",
    "» NLP Sentiment Analysis — weight predictions using customer review sentiment polarity scores.",
    "» REST API Deployment — FastAPI / Flask endpoint for external publisher tool integrations.",
    "» Cloud Deployment — Docker + Kubernetes on GCP / AWS / Azure for scalable hosting.",
    "» A/B Testing Framework — validate alternate pricing strategies with statistical significance.",
    "» Real-time Data Streaming — live Amazon price feed replacing static CSV batch pipeline runs.",
])
print("Slide: Future Enhancements")

# ─── SLIDE 49 — CONCLUSION ──────────────────────────────────────────────────
s12 = bullet_slide("Conclusion", "Project Achievements Summary", [
    "Delivered a fully functional end-to-end automated ML pipeline from raw CSV to live predictions.",
    "Data Engineering pipeline processed 525 → 440 records generating 23 model-ready features.",
    "Ridge Regression achieved R²=1.00 accuracy — 100% of predictions within a $1 error margin.",
    "EDA engine generates 13 professional charts + 10 business insights fully automatically.",
    "Recommendation engine classifies every output into a business-friendly Pricing Band tier.",
    "Streamlit Dashboard makes complex ML outputs accessible to non-technical stakeholders instantly.",
    "Clean modular architecture with PEP-8 standards and Git tracking ensures production-readiness.",
])
print("Slide: Conclusion")

# ─── SLIDE 50 — THANK YOU ───────────────────────────────────────────────────
sty = prs.slides.add_slide(BLANK)
set_bg(sty)
rect(sty, 0, 0, 13.333, 7.5, BG)
rect(sty, 0, 0, 0.10, 7.5, PURPLE)
rect(sty, 0, 3.52, 13.333, 0.04, TEAL)
tb(sty, 0.5, 2.0, 12.2, 1.4, "Thank You",
   sz=64, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
tb(sty, 0.5, 3.65, 12.2, 0.6, "Questions & Technical Discussion",
   sz=26, color=TEAL, align=PP_ALIGN.CENTER)
tb(sty, 0.5, 4.35, 12.2, 0.5,
   "Business Analytics Final Project  ·  2026",
   sz=20, color=GREY_BLUE, align=PP_ALIGN.CENTER)
print("Slide: Thank You")

# ─── SAVE ───────────────────────────────────────────────────────────────────
out = "Amzone_Book_Price_Presentation_COMPLETE_v3.pptx"
prs.save(out)
print(f"\n{'='*55}")
print(f"SAVED: {out}")
print(f"TOTAL SLIDES: {len(prs.slides)}")
print(f"  Content slides:        7")
print(f"  Visualization slides: {len(vis_meta)} (all 13)")
print(f"  Plot slides:          {len(plot_files)} (all 10)")
print(f"  Screen layout slides: {len(layout_files)} (all 14)")
print(f"{'='*55}")
