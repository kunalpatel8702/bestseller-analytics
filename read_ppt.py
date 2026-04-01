import sys
sys.stdout.reconfigure(encoding='utf-8')
from pptx import Presentation

try:
    prs = Presentation(r'C:\Doc Ref\MS24147_Final_Presentation.pptx')
    slides_data = []
    for i, slide in enumerate(prs.slides):
        slide_info = {'slide_num': i+1, 'text': []}
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = para.text.strip()
                    if t:
                        slide_info['text'].append(t)
        slides_data.append(slide_info)
    
    for s in slides_data:
        print(f"--- Slide {s['slide_num']} ---")
        for t in s['text']:
            print(f"  {t}")
except Exception as e:
    print(f'Error: {e}')
