from pptx import Presentation

try:
    prs = Presentation(r"C:\Doc Ref\MS24147_Final_Presentation.pptx")
    # Delete all existing slides
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    for slide in slides:
        xml_slides.remove(slide)
        
    # Add a blank slide using Layout 0 (Title) or 1 (Title/Content)
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Test Theme Slide"
    body = slide.placeholders[1]
    body.text = "If this works, the theme is preserved."
    
    prs.save("Test_Theme.pptx")
    print("Test_Theme.pptx saved successfully.")

except Exception as e:
    print(f"Error: {e}")
